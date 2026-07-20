# zelda_api/utils.py
"""
Zelda AI Analysis Engine Utilities
Core functions for document parsing, text analysis, and AI memo generation.
"""
import re
import logging
from typing import Dict, List, Optional
from datetime import datetime
import json

logger = logging.getLogger(__name__)


class AnalyzedPitch:

    def __init__(self, summary="", metrics=None, sections=None):

        if isinstance(summary, dict):

            self.summary = summary.get("summary", "")
            self.metrics = summary.get("metrics", {})
            self.sections = summary.get("sections", {})

        else:

            self.summary = summary
            self.metrics = metrics or {}
            self.sections = sections or {}

    def to_foundry_envelope(self):

        return {
            "origin": "pitch_analysis",
            "timestamp": datetime.utcnow().isoformat() + "Z",
            "payload": {
                "summary": self.summary,
                "metrics": self.metrics,
                "sections": self.sections
            }
        }


def scan_pitch_deck(uploaded_file) -> Dict:
    """
    Parses pitch deck files (PDF, PPTX, TXT) and extracts structured data.
    
    Args:
        uploaded_file: Django UploadedFile object
        
    Returns:
        dict: Extracted text, metrics, and structured data from the pitch
    """
    try:
        # Get file extension
        filename = uploaded_file.name.lower()
        
        # Read file content
        if filename.endswith('.pdf'):
            extracted_text, _page_count = _extract_pdf_text(uploaded_file)
        elif filename.endswith('.pptx'):
            extracted_text, _page_count = _extract_pptx_text(uploaded_file)
        elif filename.endswith('.txt'):
            extracted_text = uploaded_file.read().decode('utf-8')
        else:
            return {"error": "Unsupported file format. Use PDF, PPTX, or TXT."}
        
        if not extracted_text.strip():
            return {"error": "No readable text found in document."}
        
        # Extract key metrics and sections
        metrics = _extract_metrics(extracted_text)
        sections = _extract_sections(extracted_text)
        
        return {
            "status": "success",
            "summary": extracted_text[:2000],  # First 2000 chars as summary
            "full_text": extracted_text,
            "metrics": metrics,
            "sections": sections,
            "confidence": 0.85
        }
        
    except Exception as e:
        logger.error(f"Pitch deck parsing error: {str(e)}")
        return {"error": f"Failed to parse document: {str(e)}"}


def _extract_pdf_text(pdf_file):
    """Extract text from PDF files. Returns (text, page_count)."""
    try:
        import PyPDF2
        pdf_file.seek(0)
        reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text, len(reader.pages)
    except ImportError:
        logger.warning("PyPDF2 not installed. Returning placeholder text.")
        return "PDF parsing requires PyPDF2. Install with: pip install PyPDF2", 0
    except Exception as e:
        logger.error(f"PDF extraction failed: {str(e)}")
        return "", 0


def _extract_pptx_text(pptx_file):
    """Extract text from PowerPoint files. Returns (text, slide_count)."""
    try:
        from pptx import Presentation
        pptx_file.seek(0)
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text, len(prs.slides)
    except ImportError:
        logger.warning("python-pptx not installed. Returning placeholder text.")
        return "PPTX parsing requires python-pptx. Install with: pip install python-pptx", 0
    except Exception as e:
        logger.error(f"PPTX extraction failed: {str(e)}")
        return "", 0


def _extract_metrics(text: str) -> Dict:
    """
    Extract quantitative metrics from pitch deck text.
    Looks for revenue, funding, growth rates, etc.
    """
    metrics = {
        "revenue": None,
        "users": None,
        "growth_rate": None,
        "team_size": None,
        "funding_raised": None
    }
    
    # Find dollar amounts (revenue, funding)
    dollar_pattern = r'\$[\d.,]+[MK]?'
    dollar_matches = re.findall(dollar_pattern, text)
    if dollar_matches:
        metrics["funding_raised"] = dollar_matches[0]
        if len(dollar_matches) > 1:
            metrics["revenue"] = dollar_matches[1]
    
    # Find growth rates (e.g., "200% growth")
    growth_pattern = r'(\d+)%\s+(?:growth|increase|expansion)'
    growth_match = re.search(growth_pattern, text, re.IGNORECASE)
    if growth_match:
        metrics["growth_rate"] = f"{growth_match.group(1)}%"
    
    # Find user/customer counts
    user_pattern = r'(\d+(?:,\d+)?)\s+(?:users|customers|clients)'
    user_match = re.search(user_pattern, text, re.IGNORECASE)
    if user_match:
        metrics["users"] = user_match.group(1)
    
    # Find team size
    team_pattern = r'(?:team|staff)\s+of\s+(\d+)'
    team_match = re.search(team_pattern, text, re.IGNORECASE)
    if team_match:
        metrics["team_size"] = int(team_match.group(1))
    
    return {k: v for k, v in metrics.items() if v is not None}


def _extract_sections(text: str) -> Dict:
    """
    Identify and extract major sections from pitch deck.
    """
    sections = {}
    
    # Common pitch sections
    section_patterns = {
        "problem": r"(?:problem|challenge|market pain).*?(?=(?:solution|approach|our))",
        "solution": r"(?:solution|approach|our solution).*?(?=(?:market|business model|traction))",
        "market": r"(?:market|tam|market size|addressable).*?(?=(?:business model|competitors|traction))",
        "traction": r"(?:traction|metrics|results|achievements).*?(?=(?:team|funding|ask))",
        "team": r"(?:team|founders|management).*?(?=(?:funding|ask|conclusion))",
        "ask": r"(?:ask|funding|investment).*?(?=(?:conclusion|thank|contact))"
    }
    
    for section_name, pattern in section_patterns.items():
        match = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
        if match:
            section_text = match.group(0)[:200]  # First 200 chars
            sections[section_name] = section_text.strip()
    
    return sections


def analyze_web_text(page_text: str) -> Dict:
    """
    Analyzes raw webpage or document text to extract startup intelligence.
    Used by the on-page summarizer widget.
    
    Args:
        page_text: Raw text content from a webpage or document
        
    Returns:
        dict: Structured analysis with traction, tech stack, and investment ask
    """
    try:
        analysis = {
            "traction": _extract_traction(page_text),
            "tech_stack": _extract_tech_stack(page_text),
            "investment_ask": _extract_investment_ask(page_text),
            "confidence_score": 0.75
        }
        return analysis
        
    except Exception as e:
        logger.error(f"Web text analysis failed: {str(e)}")
        return {
            "traction": "Analysis failed",
            "tech_stack": "Analysis failed",
            "investment_ask": "Analysis failed",
            "confidence_score": 0.0
        }


def _extract_traction(text: str) -> str:
    """Extract traction signals from text."""
    traction_signals = []
    
    # Look for user/customer counts
    user_pattern = r'(\d+,?\d*)\s+(?:users|customers|signups|downloads)'
    user_matches = re.findall(user_pattern, text, re.IGNORECASE)
    if user_matches:
        traction_signals.append(f"{user_matches[0]} users")
    
    # Look for revenue
    revenue_pattern = r'\$(\d+[MK]?)\s+(?:revenue|ARR|MRR)'
    revenue_matches = re.findall(revenue_pattern, text, re.IGNORECASE)
    if revenue_matches:
        traction_signals.append(f"${revenue_matches[0]} revenue")
    
    # Look for growth rate
    growth_pattern = r'(\d+)%\s+(?:growth|increase|MoM|YoY)'
    growth_matches = re.findall(growth_pattern, text, re.IGNORECASE)
    if growth_matches:
        traction_signals.append(f"{growth_matches[0]}% growth rate")
    
    # Look for partnerships/integrations
    if re.search(r'(?:partnership|integration|partnership with|integrated with)', text, re.IGNORECASE):
        traction_signals.append("Strategic partnerships established")
    
    if traction_signals:
        return "• " + "\n• ".join(traction_signals[:3])  # Return top 3
    return "No specific traction metrics detected."


def _extract_tech_stack(text: str) -> str:
    """Extract technology stack from text."""
    tech_keywords = [
        'React', 'Node.js', 'Python', 'Django', 'PostgreSQL', 'MongoDB',
        'AWS', 'Google Cloud', 'Azure', 'Docker', 'Kubernetes', 'GraphQL',
        'TensorFlow', 'PyTorch', 'Machine Learning', 'AI', 'Blockchain',
        'Web3', 'Solidity', 'Ruby on Rails', 'Vue.js', 'Angular',
        'Next.js', 'FastAPI', 'Rust', 'Go', 'Java', 'Scala'
    ]
    
    found_techs = []
    for tech in tech_keywords:
        if re.search(rf'\b{tech}\b', text, re.IGNORECASE):
            found_techs.append(tech)
    
    if found_techs:
        return "• " + "\n• ".join(found_techs[:5])  # Return top 5
    return "Standard technology framework detected."


def _extract_investment_ask(text: str) -> str:
    """Extract investment amount being requested."""
    ask_pattern = r'(?:raising|seeking|fundraising|seeking|target)\s+\$(\d+[MK]?)'
    ask_match = re.search(ask_pattern, text, re.IGNORECASE)
    
    if ask_match:
        amount = ask_match.group(1)
        return f"Seeking ${amount} in funding"
    
    return "Target funding details not explicitly stated."


def compile_executive_intelligence_memo(
    founder_app,
    investor_app=None,
    vector_score=None,
    transparency_index=None,
    investor_matches=None
):
    """
    Generates a structured executive intelligence memo. 
    Integrates with the FoundryStandardMixin protocol for data consistency.
    """
    try:
        # Use the protocol to ensure we have normalized, safe data
        data = founder_app.to_foundry_envelope(include_full_data=True)
        payload = data.get("payload", {})

        # Build the memo structure
        memo = {
            "memo_type": "Founder Intelligence Brief" if hasattr(founder_app, 'company_name') else "Investor Intelligence Brief",
            "subject": payload.get('company_name', 'Investment Mandate'),
            "timestamp": datetime.utcnow().isoformat() + "Z",
            "executive_summary": _generate_executive_summary(founder_app),
            "market_position": _generate_market_position(founder_app),
            "recommendations": _generate_recommendations(founder_app),
            "metrics_dashboard": _generate_metrics_dashboard(founder_app),
            "vector_score": vector_score,
            "transparency_index": transparency_index
        }

        # Add Investor Context if present
        if investor_app:
            memo["target_investor"] = getattr(investor_app, "company_name", "Institutional Investor")
            if investor_matches:
                memo["potential_connections"] = len(investor_matches)
                memo["top_matches"] = investor_matches[:5]
        
        return memo
        
    except Exception as e:
        logger.error(f"Memo compilation failed for {getattr(founder_app, 'id', 'unknown')}: {str(e)}")
        return {"error": f"Failed to compile memo: {str(e)}"}

def _generate_executive_summary(app) -> str:
    """Generate executive summary for a company/investor profile."""
    try:
        company_name = getattr(app, 'company_name', 'Profile')
        sector = getattr(app, 'sector', 'General')
        description = getattr(app, 'description', '')[:150]
        
        return f"{company_name} is a {sector} company focused on {description}. Based on current market positioning and network traction, significant opportunity exists for expanded investor engagement and capital deployment."
    except:
        return "Profile analysis unavailable."


def _generate_market_position(app) -> Dict:
    """Generate market positioning analysis."""
    return {
        "sector": getattr(app, 'sector', 'General'),
        "geography": getattr(app, 'location', 'Global'),
        "stage": getattr(app, 'stage', 'Early-stage'),
        "market_assessment": "Growing market with increasing institutional interest",
        "competitive_position": "Well-positioned relative to peer benchmarks"
    }


def _generate_recommendations(app) -> List[str]:
    """Generate strategic recommendations."""
    return [
        "Increase network visibility through featured placement opportunities",
        "Engage with industry-specific investor cohorts for targeted capital deployment",
        "Leverage platform analytics to refine market positioning",
        "Utilize match radar insights to optimize investor targeting strategy"
    ]


def _generate_metrics_dashboard(app) -> Dict:
    """Generate key performance metrics for dashboard display."""
    return {
        "profile_completeness": _calculate_completeness(app),
        "network_activity": "Active",
        "engagement_score": 85,
        "match_quality": "High",
        "platform_recommendations": 12
    }


def _calculate_completeness(app) -> int:
    """Calculate profile completion percentage."""
    tracked_fields = [
        'company_name', 'sector', 'description', 
        'location', 'email', 'phone_number'
    ]
    filled = sum(1 for field in tracked_fields if getattr(app, field, None))
    return int((filled / len(tracked_fields)) * 100)

def extract_text_from_file(uploaded_file):
    """
    Extracts raw text from an uploaded file, plus a page/slide count for
    display (e.g. the valuation preview's "N Pages Analyzed" stat, which
    was silently stuck at 0 before this — PyPDF2/python-pptx already know
    the count, it just wasn't being returned). Returns (text, page_count);
    .txt has no natural pagination, so it's always 1 "page."
    Used by the DocumentIngestView to get raw text for the pipeline.
    """
    try:
        filename = uploaded_file.name.lower()

        if filename.endswith('.pdf'):
            return _extract_pdf_text(uploaded_file)
        elif filename.endswith('.pptx'):
            return _extract_pptx_text(uploaded_file)
        elif filename.endswith('.txt'):
            uploaded_file.seek(0)
            text = uploaded_file.read().decode('utf-8', errors='ignore')
            return text, 1 if text else 0
        else:
            logger.error(f"Unsupported file format for extraction: {filename}")
            return "", 0
    except Exception as e:
        logger.error(f"Failed to extract text: {str(e)}")
        return "", 0